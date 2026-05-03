import json
import math
import shutil
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scratch" / "python_deps"))

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
WORKSPACE = Path(__file__).resolve().parents[1]
ASSETS = WORKSPACE / "scratch" / "assets" / "flu-defense"
TMP = WORKSPACE / "scratch" / "office_assets"
OUT = WORKSPACE / "output" / "office_compatible.pptx"
FINAL = ROOT / "基于深度学习和多元数据的流感爆发趋势预测_答辩PPT_Office兼容版.pptx"
ORIGINAL = ROOT / "基于深度学习和多元数据的流感爆发趋势预测_答辩PPT.pptx"
BACKUP = ROOT / "基于深度学习和多元数据的流感爆发趋势预测_答辩PPT_原导出版备份.pptx"
TMP.mkdir(parents=True, exist_ok=True)
OUT.parent.mkdir(parents=True, exist_ok=True)


SLIDE_W, SLIDE_H = 13.333333, 7.5

C = {
    "canvas": "F8FAFC",
    "ink": "111827",
    "muted": "5B6472",
    "faint": "E5E7EB",
    "line": "CBD5E1",
    "teal": "0E7C7B",
    "teal2": "0A5D60",
    "blue": "2563EB",
    "sky": "38A3E8",
    "coral": "E85D4A",
    "amber": "F5B841",
    "green": "2E9E62",
    "purple": "6B46C1",
    "white": "FFFFFF",
    "black": "0B1220",
    "table": "EAF7F6",
    "row": "F1F5F9",
}

FONT = "Microsoft YaHei"
MONO = "Consolas"


def rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def emu(value):
    return Inches(value)


def set_run_font(run, size=18, color=C["ink"], bold=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold


def add_text(slide, text, x, y, w, h, size=18, color=C["ink"], bold=False,
             align="left", valign="top", font=FONT, fill=None):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(fill)
        box.line.color.rgb = rgb(fill)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = str(text).split("\n")
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = line
        set_run_font(run, size=size, color=color, bold=bold, font=font)
    return box


def add_rect(slide, x, y, w, h, fill=C["white"], line=None, radius=False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, emu(x), emu(y), emu(w), emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    return shape


def add_line(slide, x1, y1, x2, y2, color=C["line"], width=1.2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(x1), emu(y1), emu(x2), emu(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(C["canvas"])


def header(slide, title, subtitle, page):
    add_text(slide, title, 0.72, 0.42, 9.3, 0.45, 27, C["ink"], True)
    if subtitle:
        add_text(slide, subtitle, 0.74, 0.96, 9.3, 0.25, 13, C["muted"])
    add_text(slide, f"{page:02d}", 12.0, 0.5, 0.55, 0.25, 13, C["teal"], True, "right")
    add_line(slide, 0.68, 1.25, 12.68, 1.25, C["faint"], 0.8)


def footer(slide, source):
    add_line(slide, 0.68, 7.02, 12.68, 7.02, C["faint"], 0.7)
    add_text(slide, source, 0.72, 7.12, 7.9, 0.16, 7, "7B8493")
    add_text(slide, "中国海洋大学 · 计算机科学与技术本科毕业设计", 9.0, 7.12, 3.7, 0.16, 7, "7B8493", align="right")


def add_metric(slide, label, value, x, y, color=C["teal"], suffix=""):
    add_text(slide, value, x, y, 1.45, 0.45, 29, color, True, font=MONO)
    if suffix:
        add_text(slide, suffix, x + 1.35, y + 0.17, 0.5, 0.18, 13, color, True)
    add_text(slide, label, x, y + 0.54, 1.85, 0.24, 10, C["muted"])


def bullet(slide, title, body, x, y, accent=C["teal"], w=4.8):
    add_rect(slide, x, y + 0.06, 0.08, 0.42, accent)
    add_text(slide, title, x + 0.26, y, w, 0.25, 17, C["ink"], True)
    add_text(slide, body, x + 0.26, y + 0.32, w, 0.42, 12.2, C["muted"])


def pill(slide, text, x, y, w, color):
    add_rect(slide, x, y, w, 0.32, color, radius=True)
    add_text(slide, text, x, y + 0.04, w, 0.21, 10.5, C["white"], True, "center", "middle")


def image_cover(src, name, aspect):
    out = TMP / name
    img = Image.open(src).convert("RGB")
    w = 2000
    h = max(1, int(w / aspect))
    fit = ImageOps.fit(img, (w, h), method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))
    fit.save(out, quality=92)
    return out


def add_picture(slide, src, x, y, w, h=None, fit="contain", name="img.jpg"):
    src = Path(src)
    if fit == "cover":
        temp = image_cover(src, name, w / h)
        return slide.shapes.add_picture(str(temp), emu(x), emu(y), emu(w), emu(h))
    img = Image.open(src)
    aspect = img.width / img.height
    if h is None:
        h = w / aspect
    box_aspect = w / h
    if aspect > box_aspect:
        real_w = w
        real_h = w / aspect
        real_x = x
        real_y = y + (h - real_h) / 2
    else:
        real_h = h
        real_w = h * aspect
        real_x = x + (w - real_w) / 2
        real_y = y
    return slide.shapes.add_picture(str(src), emu(real_x), emu(real_y), emu(real_w), emu(real_h))


def table(slide, rows, x, y, w, row_h, col_ws, header_fill=C["table"], accent=C["teal"]):
    for r, row in enumerate(rows):
        cy = y + r * row_h
        fill = header_fill if r == 0 else (C["white"] if r % 2 else C["row"])
        add_rect(slide, x, cy, w, row_h, fill, C["line"])
        cx = x
        for c, value in enumerate(row):
            if c > 0:
                add_line(slide, cx, cy, cx, cy + row_h, "D7E2E8", 0.6)
            add_text(
                slide, value, cx + 0.08, cy + 0.11, col_ws[c] - 0.13, row_h - 0.14,
                10.5 if r else 11.2,
                accent if r == 0 else (C["ink"] if c == 0 else C["muted"]),
                bold=(r == 0 or c == 0),
                valign="middle",
            )
            cx += col_ws[c]


def bar_chart(slide, cats, vals, x, y, w, h, color=C["teal"], ymin=0, ymax=None, decimals=3):
    if ymax is None:
        ymax = max(vals) * 1.18 if vals else 1
    plot_x, plot_y = x + 0.45, y + 0.22
    plot_w, plot_h = w - 0.65, h - 0.65
    for i in range(5):
        yy = plot_y + plot_h * i / 4
        add_line(slide, plot_x, yy, plot_x + plot_w, yy, "E3EAF1", 0.6)
        val = ymax - (ymax - ymin) * i / 4
        add_text(slide, f"{val:.1f}", x, yy - 0.07, 0.36, 0.12, 6.8, C["muted"], align="right")
    zero_y = plot_y + plot_h * (ymax - 0) / (ymax - ymin) if ymin < 0 < ymax else plot_y + plot_h
    add_line(slide, plot_x, zero_y, plot_x + plot_w, zero_y, C["line"], 0.9)
    bw = plot_w / len(vals) * 0.45
    for i, (cat, val) in enumerate(zip(cats, vals)):
        cx = plot_x + plot_w * (i + 0.5) / len(vals)
        if val >= 0:
            top = plot_y + plot_h * (ymax - val) / (ymax - ymin)
            bh = zero_y - top
            label_y = top - 0.18
        else:
            top = zero_y
            bh = plot_y + plot_h * (ymax - val) / (ymax - ymin) - zero_y
            label_y = top + bh + 0.03
        add_rect(slide, cx - bw / 2, top, bw, max(0.02, bh), color)
        add_text(slide, f"{val:.{decimals}f}", cx - 0.28, label_y, 0.56, 0.14, 7.8, C["ink"], True, "center")
        add_text(slide, cat, cx - 0.5, y + h - 0.28, 1.0, 0.16, 6.9, C["ink"], align="center")


def line_chart(slide, cats, vals, x, y, w, h, color=C["teal"], ymin=0, ymax=1):
    plot_x, plot_y = x + 0.45, y + 0.28
    plot_w, plot_h = w - 0.7, h - 0.75
    for i in range(5):
        yy = plot_y + plot_h * i / 4
        add_line(slide, plot_x, yy, plot_x + plot_w, yy, "E3EAF1", 0.6)
        val = ymax - (ymax - ymin) * i / 4
        add_text(slide, f"{val:.1f}", x + 0.02, yy - 0.06, 0.33, 0.12, 6.5, C["muted"], align="right")
    pts = []
    for i, val in enumerate(vals):
        px = plot_x + plot_w * i / (len(vals) - 1)
        py = plot_y + plot_h * (ymax - val) / (ymax - ymin)
        pts.append((px, py))
        add_text(slide, cats[i], px - 0.16, y + h - 0.25, 0.32, 0.12, 7, C["muted"], align="center")
        add_text(slide, f"{val:.3f}", px - 0.24, py - 0.19, 0.48, 0.12, 7.2, C["ink"], True, "center")
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, emu(px - 0.035), emu(py - 0.035), emu(0.07), emu(0.07))
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(color)
        dot.line.fill.background()
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        add_line(slide, x1, y1, x2, y2, color, 1.8)


def make_deck():
    prs = Presentation()
    prs.slide_width = emu(SLIDE_W)
    prs.slide_height = emu(SLIDE_H)
    blank = prs.slide_layouts[6]

    def new_slide():
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        return slide

    # 1
    s = new_slide()
    add_picture(s, ASSETS / "ili_trend_cover.png", 0, 3.9, SLIDE_W, 2.25, "cover", "cover_trend.jpg")
    add_rect(s, 0, 6.05, SLIDE_W, 1.45, C["canvas"])
    add_text(s, "基于深度学习和多元数据的", 0.9, 0.92, 5.7, 0.35, 23, C["teal2"], True)
    add_text(s, "流感爆发趋势预测", 0.86, 1.46, 6.5, 0.64, 43, C["ink"], True)
    add_line(s, 0.88, 2.38, 3.75, 2.38, C["coral"], 3)
    add_text(s, "本科毕业设计答辩", 0.9, 2.7, 3.2, 0.28, 18, C["muted"], True)
    add_text(s, "国家流感中心北方省份周度 ILI% · 气象数据 · 百度指数 · iTransformer", 0.9, 3.08, 6.4, 0.24, 11.5, C["muted"])
    add_text(s, "16周输入\n预测未来4周", 9.6, 1.18, 2.0, 0.6, 19, C["teal"], True, "center")
    add_line(s, 9.42, 2.14, 11.58, 2.14, C["teal"], 4)
    add_text(s, "真实公开数据\n严格时间切分\n可复现工程链路", 9.9, 2.42, 1.7, 0.72, 16, C["ink"], True, "center")
    add_text(s, "学生：张宇鑫  |  学号：22090032057  |  指导教师：刘艳艳", 0.9, 6.32, 5.8, 0.22, 11.5, C["ink"], True)
    add_text(s, "中国海洋大学 信息科学与工程学部\n计算机科学与技术 2022级 · 2026年5月", 0.9, 6.68, 5.2, 0.4, 10, C["muted"])

    # 2
    s = new_slide(); header(s, "研究背景与目标", "从被动监测走向提前 1-4 周的趋势预警", 2)
    add_metric(s, "预测窗口", "4", 1.05, 1.78, C["coral"], "周")
    add_metric(s, "输入历史", "16", 3.35, 1.78, C["teal"], "周")
    add_metric(s, "建模特征", "29", 5.65, 1.78, C["blue"], "个")
    bullet(s, "现实问题", "流感具有季节性与非线性高峰，单一监测序列往往存在滞后，难以及时支撑医疗资源调度。", 0.95, 3.28, C["coral"], 5.2)
    bullet(s, "研究目标", "融合流感监测、气象环境和搜索行为数据，构建可复现的周度趋势预测系统。", 0.95, 4.32, C["teal"], 5.2)
    bullet(s, "答辩重点", "说明数据可信、模型合理、实验可复现，以及结论边界不外推到单一城市病例数。", 0.95, 5.36, C["blue"], 5.2)
    add_rect(s, 8.0, 1.78, 4.1, 3.65, "EAF7F6", "C6E3E1", True)
    add_text(s, "核心问题", 8.42, 2.22, 1.8, 0.28, 18, C["teal2"], True)
    add_text(s, "如何在真实周度公共卫生数据上，利用外部环境与行为信号，提前判断未来流感活动变化？", 8.42, 2.75, 3.2, 1.2, 19, C["ink"], True)
    add_line(s, 8.42, 4.18, 11.0, 4.18, C["teal"], 3)
    add_text(s, "输出形式：预测曲线、模型指标、消融解释、Web 展示原型。", 8.42, 4.45, 3.1, 0.54, 12.5, C["muted"])
    footer(s, "资料来源：项目 README、毕业论文修订核验版、开题报告与 results/reports 实验报告。")

    # 3
    s = new_slide(); header(s, "数据来源与研究口径", "正式实验使用国家流感中心北方省份真实周度 ILI% 序列", 3)
    table(s, [
        ["数据类型", "来源 / 口径", "规模", "建模作用"],
        ["流感监测", "中国国家流感中心周报，北方省份 ILI%", "849周；2010-01-04 至 2026-04-13", "预测目标 ili_rate"],
        ["气象环境", "Open-Meteo Archive，8个北方代表城市", "5594天；2011-01-01 至 2026-04-25", "温度、湿度、风速、气压"],
        ["搜索行为", "百度指数，流感/感冒/发烧关键词", "5593天；2011-01-01 至 2026-04-24", "公众关注度外生变量"],
        ["周度融合", "按 ISO 周对齐并特征工程", "797周融合；793周进入切分", "训练、验证、测试"],
    ], 0.85, 1.65, 9.45, 0.5, [1.2, 3.05, 2.55, 2.65])
    circle = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, emu(10.6), emu(1.72), emu(1.95), emu(1.95))
    circle.fill.solid(); circle.fill.fore_color.rgb = rgb("EAF7F6"); circle.line.color.rgb = rgb("BEE1DF")
    add_text(s, "北方省份\n\n公开可追溯\n周度监测", 10.82, 2.17, 1.5, 0.95, 18, C["ink"], True, "center", "middle")
    pill(s, "不再使用平谷代理", 10.6, 4.1, 1.95, C["coral"])
    pill(s, "不使用模拟目标序列", 10.6, 4.55, 1.95, C["teal"])
    pill(s, "不外推为病例数", 10.6, 5.0, 1.95, C["blue"])
    footer(s, "资料来源：README、final_data_audit.md、毕业论文第2章。")

    # 4
    s = new_slide(); header(s, "数据处理与质量审计", "从原始采集留痕到严格时序切分，降低数据泄漏风险", 4)
    add_picture(s, ASSETS / "thesis_data_pipeline.png", 0.78, 1.65, 8.1, 3.35)
    add_metric(s, "ili_rate 剩余缺失", "0", 9.7, 1.85, C["teal"], "周")
    add_metric(s, "有来源补齐", "3", 9.7, 3.02, C["coral"], "周")
    add_metric(s, "positive_rate 审计补齐", "108", 9.7, 4.18, C["blue"], "周")
    bullet(s, "质量策略", "source_manifest.json 记录来源，data_quality_report.json 复核行数、缺失、重复日期与时间边界。", 0.95, 5.46, C["teal"], 5.6)
    bullet(s, "边界说明", "positive_rate 仅作为原始监测留痕字段，默认不进入训练特征，避免不稳定 PDF 解析影响模型。", 6.9, 5.46, C["coral"], 5.4)
    footer(s, "资料来源：final_data_audit.md、data_quality_report.json、论文第2章。")

    # 5
    s = new_slide(); header(s, "预测任务与特征设计", "用过去 16 周多源变量预测未来 4 周 ILI%", 5)
    add_rect(s, 1.05, 2.03, 4.55, 1.05, "EAF2FF", "BCD4FF", True)
    add_text(s, "输入窗口", 1.38, 2.27, 1.25, 0.25, 15, C["blue"], True)
    add_text(s, "过去 16 周 × 29 个特征", 1.38, 2.62, 3.3, 0.3, 19, C["ink"], True)
    add_line(s, 5.6, 2.55, 7.62, 2.55, C["teal"], 3.5)
    tri = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE, emu(7.35), emu(2.38), emu(0.34), emu(0.34))
    tri.fill.solid(); tri.fill.fore_color.rgb = rgb(C["teal"]); tri.line.fill.background()
    add_rect(s, 7.9, 2.03, 4.55, 1.05, "FFF3E8", "FFD4B6", True)
    add_text(s, "输出目标", 8.23, 2.27, 1.25, 0.25, 15, C["coral"], True)
    add_text(s, "未来 1-4 周 ILI%", 8.23, 2.62, 3.2, 0.3, 19, C["ink"], True)
    table(s, [
        ["特征组", "代表变量", "设计目的"],
        ["流感历史", "ili_rate、1/2/4周滞后、4/8周滚动统计", "捕捉自相关与季节高峰"],
        ["气象环境", "温度、湿度、风速、气压、体感/交互项", "引入环境变化背景"],
        ["搜索行为", "流感、感冒、发烧搜索指数及变化率", "反映公众关注度变化"],
        ["时间特征", "week/month 周期编码、流感季标记", "表达周期性和季节性"],
    ], 1.05, 3.68, 11.2, 0.46, [1.55, 4.85, 4.8], "F1F5F9", C["blue"])
    footer(s, "资料来源：experiment_summary.json、config/config.yaml、论文第2.3节。")

    # 6
    s = new_slide(); header(s, "模型与系统设计", "iTransformer 建模跨变量关系，工程链路覆盖采集、训练、评估和展示", 6)
    add_text(s, "iTransformer 的关键变化", 0.98, 1.68, 4.4, 0.32, 17, C["teal2"], True)
    bullet(s, "变量作为 token", "将每个变量的历史窗口序列作为一个 token，在变量维度计算注意力。", 0.95, 2.16, C["teal"], 4.9)
    bullet(s, "目标 token 输出", "学习流感历史、气象与搜索行为之间的关联后，映射为未来 4 周预测。", 0.95, 3.22, C["blue"], 4.9)
    bullet(s, "基准模型对照", "设置 LSTM、DLinear、ARIMA，分别代表循环网络、线性分解与传统统计方法。", 0.95, 4.28, C["coral"], 4.9)
    add_picture(s, ASSETS / "thesis_system_architecture.png", 6.0, 1.72, 6.45, 2.65)
    table(s, [
        ["层次", "项目实现"],
        ["数据层", "CSV / source_manifest.json"],
        ["处理层", "清洗、周聚合、特征工程、切分"],
        ["模型层", "iTransformer + 基准模型"],
        ["展示层", "Streamlit Web 原型"],
    ], 6.4, 4.75, 5.15, 0.39, [1.25, 3.9])
    footer(s, "资料来源：src/models/itransformer.py、scripts/train.py、web/app.py、论文第3章与第5章。")

    # 7
    s = new_slide(); header(s, "实验设计", "按时间顺序切分，避免未来信息进入训练", 7)
    x0, y0, total = 1.1, 2.48, 9.4
    train_w, val_w, test_w = 6.55, 1.42, 1.43
    add_rect(s, x0, y0, train_w, 0.6, C["teal"])
    add_rect(s, x0 + train_w, y0, val_w, 0.6, C["blue"])
    add_rect(s, x0 + train_w + val_w, y0, test_w, 0.6, C["coral"])
    add_text(s, "训练集 555行\n2011-01-31 至 2021-09-20", x0 + 0.18, y0 + 0.08, train_w - 0.36, 0.42, 12, C["white"], True, "center", "middle")
    add_text(s, "验证集 119行\n2021-09-27 至 2024-01-01", x0 + train_w + 0.04, y0 + 0.08, val_w - 0.08, 0.42, 8.4, C["white"], True, "center", "middle")
    add_text(s, "测试集 119行\n2024-01-08 至 2026-04-13", x0 + train_w + val_w + 0.04, y0 + 0.08, test_w - 0.08, 0.42, 8.4, C["white"], True, "center", "middle")
    add_line(s, x0, y0 + 0.9, x0 + total, y0 + 0.9, C["line"], 1)
    add_text(s, "2011", x0 - 0.02, y0 + 1.04, 0.55, 0.14, 8, C["muted"])
    add_text(s, "2021", x0 + train_w - 0.2, y0 + 1.04, 0.55, 0.14, 8, C["muted"])
    add_text(s, "2024", x0 + train_w + val_w - 0.2, y0 + 1.04, 0.55, 0.14, 8, C["muted"])
    add_text(s, "2026", x0 + total - 0.42, y0 + 1.04, 0.55, 0.14, 8, C["muted"], align="right")
    table(s, [
        ["训练配置", "取值"],
        ["随机种子", "42"],
        ["最大轮数 / 早停", "300 / patience=30"],
        ["优化器", "Adam，学习率 0.001，权重衰减 0.01"],
        ["评价指标", "RMSE、MAE、MAPE、R2、峰值指标"],
    ], 1.1, 4.18, 5.4, 0.4, [1.5, 3.9], "F1F5F9", C["blue"])
    bullet(s, "实验边界", "指标用于当前北方地区真实周度监测口径下的模型比较，不作为单一城市病例数预测结论。", 7.2, 4.4, C["coral"], 4.6)
    footer(s, "资料来源：experiment_summary.json、config/config.yaml、论文第4.1节。")

    # 8
    s = new_slide(); header(s, "模型对比结果", "iTransformer 在整体误差和解释方差上表现最好", 8)
    add_metric(s, "iTransformer RMSE", "0.772", 0.98, 1.82, C["teal"])
    add_metric(s, "MAE", "0.491", 3.1, 1.82, C["blue"])
    add_metric(s, "MAPE", "11.51", 5.35, 1.82, C["coral"], "%")
    add_metric(s, "R²", "0.538", 7.8, 1.82, C["green"])
    cats = ["iTransformer", "ARIMA", "DLinear", "LSTM"]
    bar_chart(s, cats, [0.772, 0.980, 1.031, 1.366], 0.95, 3.45, 5.6, 2.8, C["teal"], 0, 1.4)
    bar_chart(s, cats, [0.538, 0.450, 0.176, -0.444], 7.0, 3.45, 5.6, 2.8, C["green"], -0.6, 0.6)
    add_text(s, "结论：当前数据版本下，iTransformer 的整体 RMSE 最低；ARIMA 在短期与峰值指标上仍有竞争力。", 0.95, 6.32, 10.0, 0.25, 14, C["ink"], True)
    footer(s, "资料来源：all_metrics.json、experiment_brief.md、论文第4.2节。")

    # 9
    s = new_slide(); header(s, "预测效果与多步误差", "越远的预测步长不确定性越高，符合多步时间序列规律", 9)
    add_picture(s, ASSETS / "iTransformer_predictions.png", 0.67, 1.55, 7.75, 4.65, fit="cover", name="pred_cover.jpg")
    add_rect(s, 8.75, 1.75, 3.62, 2.05, C["white"], C["line"], True)
    add_text(s, "分 Horizon 误差", 9.1, 2.08, 2.6, 0.28, 17, C["teal2"], True)
    line_chart(s, ["H1", "H2", "H3", "H4"], [0.525, 0.676, 0.828, 0.983], 9.0, 2.45, 2.95, 1.2, C["teal"], 0, 1.0)
    add_metric(s, "H1 RMSE", "0.525", 9.05, 4.42, C["teal"])
    add_metric(s, "H4 RMSE", "0.983", 10.82, 4.42, C["coral"])
    add_text(s, "解释：第1周误差最低，预测距离增加后 RMSE 和 MAE 上升；高峰阶段仍是主要误差来源。", 9.05, 5.48, 3.1, 0.65, 13, C["ink"], True)
    footer(s, "资料来源：iTransformer_predictions.png、horizon_metrics.json、论文第4.3节。")

    # 10
    s = new_slide(); header(s, "消融实验与解释", "外生变量的贡献需要结合聚合噪声和流行季背景解释", 10)
    cats2 = ["仅流感历史", "流感+气象", "三源融合", "流感+搜索"]
    vals2 = [0.761, 0.767, 0.770, 0.788]
    bar_chart(s, cats2, vals2, 0.95, 1.95, 6.0, 3.9, C["blue"], 0, 0.85)
    table(s, [
        ["特征组合", "特征数", "RMSE", "R²"],
        ["仅流感历史", "13", "0.761", "0.552"],
        ["流感+气象", "20", "0.767", "0.544"],
        ["三源融合", "29", "0.770", "0.540"],
        ["流感+搜索", "22", "0.788", "0.519"],
    ], 7.45, 1.95, 4.6, 0.46, [1.7, 0.85, 1.0, 1.05])
    bullet(s, "实验发现", "仅流感历史 RMSE 最低，流感+气象与三源融合非常接近。", 7.45, 4.9, C["teal"], 4.6)
    bullet(s, "答辩表述", "不能简单说“搜索指数一定提升性能”；更稳妥的结论是外生变量提供辅助解释，但受代表城市聚合口径影响。", 7.45, 5.9, C["coral"], 4.6)
    footer(s, "资料来源：ablation_report.md、ablation_metrics.csv、论文第4.4节。")

    # 11
    s = new_slide(); header(s, "Web 展示与工程复现", "答辩演示可从结果文件直接读取，不依赖临时手工整理", 11)
    add_picture(s, ASSETS / "thesis_system_architecture.png", 0.85, 1.7, 6.9, 2.75)
    add_text(s, "Streamlit 原型功能", 8.25, 1.9, 3.0, 0.32, 17, C["teal2"], True)
    bullet(s, "数据总览", "多源时序、样本数量、时间范围、相关性矩阵。", 8.25, 2.48, C["teal"], 3.8)
    bullet(s, "预测分析", "训练曲线、预测曲线、误差分布、注意力热力图。", 8.25, 3.4, C["blue"], 3.8)
    bullet(s, "对比实验", "多模型指标、消融实验与预警调度演示。", 8.25, 4.32, C["coral"], 3.8)
    add_rect(s, 1.02, 5.0, 6.5, 0.92, C["black"], radius=True)
    add_text(s, "python scripts/train.py --skip-collect\npython scripts/run_ablation.py --skip-collect\nstreamlit run web/app.py", 1.25, 5.15, 5.9, 0.58, 9.5, "DFF7EF", font=MONO)
    footer(s, "资料来源：README、web/app.py、scripts/train.py、scripts/run_ablation.py。")

    # 12
    s = new_slide(); header(s, "结论与展望", "完成了真实数据、多源融合、深度模型和可视化系统的闭环", 12)
    add_text(s, "主要结论", 0.95, 1.8, 3.0, 0.34, 20, C["teal2"], True)
    bullet(s, "数据可信", "正式实验切换为国家流感中心北方省份周度 ILI%，并保留缺失补齐和来源审计。", 0.95, 2.36, C["teal"], 5.7)
    bullet(s, "模型有效", "iTransformer 在测试集上 RMSE=0.772、MAE=0.491、MAPE=11.51%、R²=0.538。", 0.95, 3.42, C["blue"], 5.7)
    bullet(s, "工程完整", "项目覆盖采集、预处理、特征工程、训练评估、消融实验和 Web 展示。", 0.95, 4.48, C["coral"], 5.7)
    add_text(s, "不足与展望", 7.45, 1.8, 3.0, 0.34, 20, C["coral"], True)
    bullet(s, "区域聚合仍可细化", "百度指数和气象采用代表城市简单聚合，后续可引入人口权重或省级分层。", 7.45, 2.36, C["coral"], 4.8)
    bullet(s, "峰值预测仍需增强", "高发期误差更明显，可尝试高峰样本加权、分季节训练和概率预测。", 7.45, 3.42, C["blue"], 4.8)
    bullet(s, "解释性不等于因果", "注意力热力图只能辅助理解信息流，不能直接作为流感传播因果结论。", 7.45, 4.48, C["teal"], 4.8)
    add_text(s, "谢谢各位老师，请批评指正", 3.0, 6.12, 7.4, 0.48, 29, C["ink"], True, "center")
    footer(s, "资料来源：毕业论文第6-7章、experiment_brief.md、ablation_report.md。")

    return prs


if __name__ == "__main__":
    prs = make_deck()
    prs.save(OUT)
    # Verify python-pptx can reopen the generated file.
    reopened = Presentation(str(OUT))
    if len(reopened.slides) != 12:
        raise RuntimeError(f"Unexpected slide count: {len(reopened.slides)}")
    if ORIGINAL.exists() and not BACKUP.exists():
        shutil.copy2(ORIGINAL, BACKUP)
    shutil.copy2(OUT, FINAL)
    shutil.copy2(OUT, ORIGINAL)
    print(json.dumps({
        "office_pptx": str(FINAL),
        "overwritten_original": str(ORIGINAL),
        "backup": str(BACKUP),
        "slide_count": len(reopened.slides),
        "size": OUT.stat().st_size,
    }, ensure_ascii=False, indent=2))
